from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import sys

try:
    print("Initializing presentation...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    def add_title_slide(prs, title, subtitle):
        print(f"Adding title slide: {title}")
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Subtitle is usually placeholder 1
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        return slide

    def add_content_slide(prs, title, content_list):
        print(f"Adding content slide: {title}")
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for item in content_list:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
        return slide

    def add_image_slide(prs, title, image_path):
        print(f"Adding image slide: {title}")
        slide_layout = prs.slide_layouts[5] # Title Only
        slide = prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        if os.path.exists(image_path):
            print(f"  Found image: {image_path}")
            left = Inches(1)
            top = Inches(1.5)
            height = Inches(5)
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            print(f"  WARNING: Image not found: {image_path}")
            # Add text saying image missing
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = f"Image not found: {image_path}"
        
        return slide

    # Slide 1: Title
    add_title_slide(prs, 
        "University E-Voting System",
        "Milestone One: Discovery & Architecture\nDecember 2025")

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        "🎯 Objective: Secure, accessible digital voting platform",
        "📱 Mobile-first architecture with React Native",
        "🔐 OTP-based authentication for voters",
        "🏛️ Role-based access (Admin, Officer, Candidate, Voter)",
        "📊 Real-time results and audit logging"
    ])

    # Slide 3: Key Features
    add_content_slide(prs, "Key Features", [
        "✅ Voter Authentication via OTP (Email)",
        "✅ Candidate Nomination & Approval Workflow",
        "✅ Secret Ballot with Single-use Tokens",
        "✅ Real-time Vote Counting & Results",
        "✅ Immutable Audit Logs",
        "✅ Admin Dashboard for Election Management"
    ])

    # Slide 4: Stakeholders
    add_content_slide(prs, "Stakeholder Mapping", [
        "👤 Administrators - System configuration & oversight",
        "👤 Returning Officers - Nomination approval & monitoring",
        "👤 Candidates - Submit nominations & manifestos",
        "👤 Voters - Verify identity & cast votes"
    ])

    # Slide 5: Voter Journey
    add_image_slide(prs, "Voter Journey", "docs/images/voter_journey.png")

    # Slide 6: Candidate Journey
    add_image_slide(prs, "Candidate Journey", "docs/images/candidate_journey.png")

    # Slide 7: Admin Journey
    add_image_slide(prs, "Admin & Officer Journey", "docs/images/admin_journey.png")

    # Slide 8: System Architecture
    add_image_slide(prs, "Entity Relationship Diagram", "docs/images/erd.png")

    # Slide 9: Technology Stack - Backend
    add_content_slide(prs, "Technology Stack: Backend", [
        "⚙️ Runtime: Node.js",
        "🚀 Framework: Express.js",
        "💾 Database: MySQL",
        "🔧 ORM: Prisma",
        "🔐 Authentication: JWT + Bcrypt",
        "📧 Email: Nodemailer"
    ])

    # Slide 10: Technology Stack - Mobile
    add_content_slide(prs, "Technology Stack: Mobile", [
        "📱 Framework: React Native (Expo)",
        "💻 Language: JavaScript",
        "🔄 State Management: Context API",
        "💾 Storage: AsyncStorage",
        "🌐 Networking: Axios"
    ])

    # Slide 11: Security Features
    add_content_slide(prs, "Security & Risk Mitigation", [
        "🔒 Role-Based Access Control (RBAC)",
        "🔑 Two-Factor Authentication (Reg No + OTP)",
        "🎫 Single-use Ballot Tokens",
        "📝 Immutable Audit Logs",
        "🔐 Encrypted Password Storage (Bcrypt)",
        "⚡ Rate Limiting on API endpoints"
    ])

    # Slide 12: CI/CD & Version Control
    print("Adding Development Workflow slide")
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Development Workflow"
    
    tf = slide.placeholders[1].text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "📦 Two Repositories:"
    p.font.size = Pt(18)
    
    # Backend Link
    p = tf.add_paragraph()
    p.level = 0
    r1 = p.add_run()
    r1.text = "   • Backend API: "
    r1.font.size = Pt(18)
    
    r2 = p.add_run()
    r2.text = "Evoting-Backend-API"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0, 0, 255)
    r2.font.underline = True
    r2.hyperlink.address = "https://github.com/JDanielZ5G/Evoting-Backend-API.git"
    
    # Mobile Link
    p = tf.add_paragraph()
    p.level = 0
    r1 = p.add_run()
    r1.text = "   • Mobile App: "
    r1.font.size = Pt(18)
    
    r2 = p.add_run()
    r2.text = "Evoting_mobile_group3"
    r2.font.size = Pt(18)
    r2.font.color.rgb = RGBColor(0, 0, 255)
    r2.font.underline = True
    r2.hyperlink.address = "https://github.com/JDanielZ5G/Evoting_mobile_group3.git"
    
    # Other items
    for item in ["🌿 Feature-branch workflow", "✅ GitHub Actions CI/CD", "📝 Conventional Commits", "🧪 Automated Testing"]:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)

    # Slide 13: Future Enhancements
    add_content_slide(prs, "Future Plans", [
        "🔐 Biometric Authentication (FaceID/TouchID)",
        "⛓️ Blockchain Integration for vote hashes",
        "📲 Push Notifications for real-time updates",
        "🌍 Multi-language Support",
        "📊 Advanced Analytics Dashboard"
    ])

    # Slide 14: Thank You
    print("Adding Thank You slide")
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank or Title Only? Layout 6 is usually Blank. Let's use Title Only (5) or Title Slide (0)
    # Actually layout 6 is often 'Blank'. Let's check.
    # Standard layouts: 0=Title, 1=Title+Content, 2=Section Header, 3=Two Content, 4=Comparison, 5=Title Only, 6=Blank, 7=Content with Caption, 8=Picture with Caption
    # Let's use Layout 0 for Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if slide.shapes.title:
        slide.shapes.title.text = "Thank You"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Questions?\n\nUniversity E-Voting System\nMilestone One Presentation"

    output_path = 'docs/MILESTONE_ONE_PRESENTATION.pptx'
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")

except Exception as e:
    print(f"ERROR: {str(e)}")
    import traceback
    traceback.print_exc()

